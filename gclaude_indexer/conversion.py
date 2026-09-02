# GClaude Indexer — document collection indexer
# Copyright (C) 2026  Alex Camacho Castilho
#
# This program is free software: you can redistribute it and/or modify it
# under the terms of the GNU General Public License as published by the Free
# Software Foundation, either version 3 of the License, or (at your option)
# any later version. See the LICENSE file for details.

"""Conversion by type, OCR, and slicing of large PDFs (section 5, steps 2 and 3).

Every file with `status = "discovered"` is processed once. The original is
never read for writing nor altered — all output (OCR'd PDF, sliced blocks,
extracted text) goes to `output_folder`. A corrupted file, or one that fails
at any step, gets `status = "failed"`, has the error recorded in `file.error`
and in an event, and processing moves on to the next file.
"""

from __future__ import annotations

import email
import os
import subprocess
import sys
from collections.abc import Callable
from concurrent.futures import CancelledError, ProcessPoolExecutor, as_completed
from concurrent.futures.process import BrokenProcessPool
from dataclasses import dataclass
from email import policy
from pathlib import Path

import fitz  # PyMuPDF

from . import no_window, tools
from .config import ProjectConfig
from .events import record_event
from .parallelism import _physical_cores, workers_for
from .file_types import category_of_extension

MIN_AVG_CHARS_PER_PAGE = 100
ERROR_MESSAGE_CHAR_LIMIT = 500


@dataclass
class ConversionResult:
    converted: int = 0
    ocr_applied: int = 0
    sliced: int = 0
    blocks_generated: int = 0
    failed: int = 0


class ConversionError(Exception):
    """Failure converting a file. The message must not contain document content."""


def _destination_path(output_dir: Path, subfolder: str, relative_path: str, new_extension: str) -> Path:
    destination = output_dir / subfolder / Path(relative_path).with_suffix(new_extension)
    destination.parent.mkdir(parents=True, exist_ok=True)
    return destination


def _pdf_needs_ocr(path: Path) -> tuple[bool, int]:
    """Returns `(needs_ocr, page_count)`. A PDF with no text layer, or with
    fewer than 100 characters per page on average, needs OCR (section 5,
    step 2)."""
    try:
        document = fitz.open(path)
    except Exception as exc:
        raise ConversionError(f"PDF ilegível ou corrompido: {exc}") from exc

    try:
        page_count = document.page_count
        if page_count == 0:
            raise ConversionError("PDF sem páginas")

        total_chars = sum(len(page.get_text()) for page in document)
        avg = total_chars / page_count
        return avg < MIN_AVG_CHARS_PER_PAGE, page_count
    finally:
        document.close()


def _ocr_environment() -> dict:
    """Environment for the OCR subprocess, with the console-window hook on
    its `PYTHONPATH`.

    Two entries, both needed and for different reasons:

    * the project root, so `python -m gclaude_indexer._ocr_runner` resolves
      at all — the worker process that builds this command may have been
      `spawn`ed with a different working directory;
    * `_nowindow_site`, whose `sitecustomize.py` the `site` module imports
      into **every** interpreter started with this environment. That is
      what reaches ocrmypdf's own `--jobs` pool workers, which are the
      direct parents of `tesseract.exe` and would otherwise each flash a
      console window (see `no_window.py`).
    """
    package_folder = Path(__file__).resolve().parent
    entries = [str(package_folder.parent), str(package_folder / "_nowindow_site")]
    existing = os.environ.get("PYTHONPATH", "")
    if existing:
        entries.append(existing)
    return {
        **os.environ,
        "PYTHONPATH": os.pathsep.join(entries),
        "PYTHONDONTWRITEBYTECODE": "1",
        # ocrmypdf finds `tesseract` and `gswin64c` through `PATH` and
        # takes no configuration for either, so the locations the installer
        # recorded have to reach it this way (see `tools.py`).
        "PATH": tools.path_with_tools(),
    }


def _run_ocrmypdf(source: Path, destination: Path, ocr_language: str, jobs: int = 1) -> None:
    if destination.exists():
        destination.unlink()

    command = [
        sys.executable,
        # `-B` as well as `PYTHONDONTWRITEBYTECODE` in the environment: the
        # variable is the one that reaches ocrmypdf's own pool workers, but
        # it can be lost if a caller rebuilds the environment, and the
        # project folder is synced by Google Drive and must never receive a
        # `__pycache__` (section 11.1). The flag cannot be lost.
        "-B",
        "-m",
        # Not `ocrmypdf` directly: `_ocr_runner` is the same command line
        # with console windows suppressed first — see its docstring and
        # `no_window.py`. Ghostscript and Tesseract are started by this
        # process, and a windowless parent is exactly what makes each of
        # them pop a black window on the desktop.
        "gclaude_indexer._ocr_runner",
        "--skip-text",
        "--language",
        ocr_language,
        "--deskew",
        "--jobs",
        str(max(1, jobs)),
        # Section 5 also asks for "--clean", but that option depends on the
        # `unpaper` binary, which has no official Windows distribution and
        # is not on the installable dependency list in section 10.3. Omitted
        # here so OCR works on this platform; revisit if `unpaper` becomes
        # installable on Windows.
        str(source),
        str(destination),
    ]
    result = subprocess.run(
        command,
        shell=False,
        capture_output=True,
        text=True,
        env=_ocr_environment(),
        startupinfo=no_window.hidden_startupinfo(),
        creationflags=(subprocess.CREATE_NO_WINDOW if sys.platform == "win32" else 0),
    )

    if result.returncode != 0:
        raise ConversionError(
            f"ocrmypdf falhou (código {result.returncode}): "
            f"{result.stderr.strip()[-ERROR_MESSAGE_CHAR_LIMIT:]}"
        )


def slice_pdf(source_path: Path, destination_dir: Path, base_name: str, pages_per_block: int) -> list[Path]:
    """Cuts `source_path` into blocks of at most `pages_per_block` pages.

    Block name: `<base>_p<start>-<end>.pdf`, preserving the real (1-based)
    page numbering of the source document in the name.
    """
    document = fitz.open(source_path)
    try:
        page_count = document.page_count
        blocks: list[Path] = []

        for start in range(0, page_count, pages_per_block):
            end = min(start + pages_per_block, page_count) - 1
            start_1based = start + 1
            end_1based = end + 1

            block = fitz.open()
            block.insert_pdf(document, from_page=start, to_page=end)
            destination_dir.mkdir(parents=True, exist_ok=True)
            block_path = destination_dir / f"{base_name}_p{start_1based}-{end_1based}.pdf"
            block.save(block_path)
            block.close()
            blocks.append(block_path)

        return blocks
    finally:
        document.close()


def _extract_text_docx(path: Path) -> str:
    from docx import Document

    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _extract_text_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    text_rows = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if value is None else str(value) for value in row]
            text_rows.append("\t".join(cells))
    workbook.close()
    return "\n".join(text_rows)


def _extract_text_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    parts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text:
                        parts.append(text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            parts.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(parts)


def _extract_text_eml(path: Path) -> str:
    with open(path, "rb") as file:
        message = email.message_from_binary_file(file, policy=policy.default)

    parts = [f"From: {message.get('From', '')}", f"Subject: {message.get('Subject', '')}"]
    body = message.get_body(preferencelist=("plain", "html"))
    if body is not None:
        parts.append(body.get_content())
    return "\n\n".join(parts)


def _extract_text_msg(path: Path) -> str:
    import extract_msg

    message = extract_msg.Message(str(path))
    try:
        parts = [f"From: {message.sender or ''}", f"Subject: {message.subject or ''}"]
        if message.body:
            parts.append(message.body)
        return "\n\n".join(parts)
    finally:
        message.close()


def _extract_text_image(path: Path, ocr_language: str) -> str:
    from PIL import Image
    import pytesseract

    with Image.open(path) as image:
        image.load()
        return pytesseract.image_to_string(image, lang=ocr_language)


def _process_file(config: ProjectConfig, row: dict, ocrmypdf_jobs: int) -> dict:
    """Converts one file and returns what the main process needs to write to
    the database — never touches `conn`.

    Runs both on the sequential path (direct call) and inside a
    `ProcessPoolExecutor` worker process (via `pickle`, on Windows with
    `spawn`): that's why it's a module-level function, takes only
    serializable arguments (`config` is a plain dataclass; `row` is a
    `dict`, not a `sqlite3.Row`/`sqlite3.Connection`), and never receives
    `conn`.
    """
    relative_path = row["relative_path"]
    name = row["name"]
    extension_with_dot = "." + row["extension"]
    category = category_of_extension(extension_with_dot)
    source_path = Path(config.source_folder) / relative_path
    output_dir = Path(config.output_folder)

    # (step, level, i18n key, params) — rendered by `_apply_result`, which
    # runs in the main process and has `conn`/`language`; this function may
    # run in a worker process (see the docstring above) and never does.
    events: list[tuple[str, str, str, dict]] = []
    ocr_applied = False
    sliced = False
    blocks_generated = 0

    if category == "pdf":
        needs_ocr, page_count = _pdf_needs_ocr(source_path)
        path_to_slice = source_path

        if needs_ocr:
            ocr_destination = _destination_path(output_dir, "converted", relative_path, ".pdf")
            _run_ocrmypdf(source_path, ocr_destination, config.ocr_language, jobs=ocrmypdf_jobs)
            path_to_slice = ocr_destination
            ocr_applied = True
            events.append((
                "conversion", "info", "log.conversion.ocr_applied",
                {"name": name, "ocr_language": config.ocr_language},
            ))

        if page_count > config.pages_per_block:
            blocks_dir = output_dir / "blocks" / Path(relative_path).parent
            blocks = slice_pdf(path_to_slice, blocks_dir, Path(relative_path).stem, config.pages_per_block)
            sliced = True
            blocks_generated = len(blocks)
            events.append((
                "conversion", "info", "log.conversion.sliced", {"name": name, "blocks": len(blocks)},
            ))

        sql = "UPDATE file SET status = 'converted', error = NULL, page_count = ?, needs_ocr = ? WHERE id = ?"
        params = (page_count, int(needs_ocr), row["id"])

    elif category == "imagens":
        text = _extract_text_image(source_path, config.ocr_language)
        destination = _destination_path(output_dir, "converted", relative_path, ".txt")
        destination.write_text(text, encoding="utf-8")
        ocr_applied = True
        events.append((
            "conversion", "info", "log.conversion.ocr_applied",
            {"name": name, "ocr_language": config.ocr_language},
        ))
        sql = "UPDATE file SET status = 'converted', error = NULL, page_count = 1, needs_ocr = 1 WHERE id = ?"
        params = (row["id"],)

    elif category in ("docx", "xlsx", "pptx", "text", "web_dados", "email"):
        extension = row["extension"]
        if category == "docx":
            text = _extract_text_docx(source_path)
        elif category == "xlsx":
            text = _extract_text_xlsx(source_path)
        elif category == "pptx":
            text = _extract_text_pptx(source_path)
        elif category in ("text", "web_dados"):
            text = source_path.read_text(encoding="utf-8", errors="replace")
        elif extension == "msg":
            text = _extract_text_msg(source_path)
        else:
            text = _extract_text_eml(source_path)

        destination = _destination_path(output_dir, "converted", relative_path, ".txt")
        destination.write_text(text, encoding="utf-8")
        sql = "UPDATE file SET status = 'converted', error = NULL, page_count = 1, needs_ocr = 0 WHERE id = ?"
        params = (row["id"],)

    else:
        # Only reached with the special "all" category selected (explicit
        # user request: index any non-binary/executable extension, even
        # without a dedicated extractor). Reads as text, best-effort — never
        # fails on a binary file with no extractor; worst case it produces a
        # low-quality result, not an error.
        text = source_path.read_text(encoding="utf-8", errors="replace")
        destination = _destination_path(output_dir, "converted", relative_path, ".txt")
        destination.write_text(text, encoding="utf-8")
        sql = "UPDATE file SET status = 'converted', error = NULL, page_count = 1, needs_ocr = 0 WHERE id = ?"
        params = (row["id"],)

    return {
        "id": row["id"],
        "name": name,
        "sql": sql,
        "params": params,
        "events": events,
        "ocr_applied": ocr_applied,
        "sliced": sliced,
        "blocks_generated": blocks_generated,
    }


def _apply_result(conn, result: ConversionResult, info: dict, language: str | None = None) -> None:
    conn.execute(info["sql"], info["params"])
    for step, level, key, params in info["events"]:
        record_event(conn, step, level, key, params, language=language)
    result.converted += 1
    result.ocr_applied += int(info["ocr_applied"])
    result.sliced += int(info["sliced"])
    result.blocks_generated += info["blocks_generated"]


def _record_failure(conn, result: ConversionResult, row, exc: Exception, language: str | None = None) -> None:
    error_message = str(exc)[:ERROR_MESSAGE_CHAR_LIMIT]
    conn.execute(
        "UPDATE file SET status = 'failed', error = ? WHERE id = ?",
        (error_message, row["id"]),
    )
    result.failed += 1
    record_event(
        conn, "conversion", "error", "log.conversion.failed",
        {"name": row["name"], "error": error_message}, language=language,
    )


def _convert_sequential(
    conn, config: ProjectConfig, rows, ocrmypdf_jobs: int, should_stop, result, language: str | None = None
) -> None:
    for row in rows:
        if should_stop is not None and should_stop():
            break

        try:
            info = _process_file(config, dict(row), ocrmypdf_jobs)
            _apply_result(conn, result, info, language=language)
        except Exception as exc:
            _record_failure(conn, result, row, exc, language=language)
        conn.commit()  # commit per file (Phase 11): the progress bar cannot stall


def _convert_in_parallel(
    conn, config: ProjectConfig, rows, workers: int, ocrmypdf_jobs: int, should_stop, result,
    language: str | None = None,
) -> None:
    """One file per worker process. `sqlite3.Connection` never crosses a
    process boundary — each worker only returns the result (a serializable
    dict); the one that always writes to the database is this main process.

    `initializer=no_window.install`: each worker is `spawn`ed as a fresh
    interpreter, so the main process's console-window suppression does not
    carry over on its own. A worker that handles an image calls
    `pytesseract`, which starts `tesseract.exe` through a `subprocess.Popen`
    of its own that we do not control — without the initializer that is one
    visible console window per image."""
    with ProcessPoolExecutor(max_workers=workers, initializer=no_window.install) as executor:
        futures: dict = {}
        for row in rows:
            if should_stop is not None and should_stop():
                break
            future = executor.submit(_process_file, config, dict(row), ocrmypdf_jobs)
            futures[future] = row

        stop_request_handled = False
        pool_broken = False
        for future in as_completed(futures):
            row = futures[future]
            try:
                info = future.result()
            except CancelledError:
                continue
            except BrokenProcessPool:
                # A worker died (e.g. `MemoryError` on a large PDF) and took
                # the whole pool down with it — ALL still-pending futures
                # receive this same exception, including files that never
                # got to run (I3, Phase 13 final review). Unlike a real
                # per-file failure: we don't mark it 'failed' (which would
                # exclude the file from every future rerun); we leave the
                # original status ('discovered') for `convert()` to resume
                # this file on the next run.
                if not pool_broken:
                    pool_broken = True
                    record_event(conn, "conversion", "error", "log.pool_broken", language=language)
                continue
            except Exception as exc:
                _record_failure(conn, result, row, exc, language=language)
                conn.commit()
                continue

            _apply_result(conn, result, info, language=language)
            conn.commit()  # commit per file (Phase 11), even in parallel

            if not stop_request_handled and should_stop is not None and should_stop():
                stop_request_handled = True
                # The pause button must never go unanswered: cancel every
                # future not yet started. The ones already running finish
                # (there's no way to interrupt a running process here), but
                # no new file starts.
                for pending in futures:
                    if not pending.done():
                        pending.cancel()


def convert(
    conn, config: ProjectConfig, should_stop: Callable[[], bool] | None = None, language: str | None = None
) -> ConversionResult:
    # Covers the sequential path, which runs `pytesseract` (and therefore
    # `tesseract.exe`) inside this very process. The parallel path installs
    # it per worker instead, through the pool's initializer.
    no_window.install()

    result = ConversionResult()

    rows = conn.execute(
        "SELECT * FROM file WHERE status = 'discovered' ORDER BY relative_path"
    ).fetchall()

    workers = workers_for(config.parallelism)
    # With 0 or 1 file there's nothing to parallelize — spawning processes
    # would cost more than processing directly (Windows' `spawn` has a real
    # startup cost), so the sequential path is used even with an
    # "automatic"/"maximum" mode configured.
    if workers <= 1 or len(rows) <= 1:
        # With no pool of processes competing for the same cores, `--jobs 1`
        # avoids the whole problem below at no cost (it's a single file).
        _convert_sequential(conn, config, rows, 1, should_stop, result, language=language)
    else:
        # Each pool worker runs its own `ocrmypdf --jobs N`. The two axes of
        # parallelism (pool processes x ocrmypdf threads) multiply: using the
        # same value for both (original Task 14 decision) gives, in "maximum"
        # mode with 8 cores, 8 processes x `--jobs 8` = up to 64 simultaneous
        # tesseracts competing for 8 cores. Task 14 measured a 6.73x gain with
        # 3-page PDFs — little RAM per file; a real scanned collection has far
        # more pages per file, and memory per tesseract scales by a different
        # order of magnitude. Combined with I3 (one worker running out of
        # memory took down the whole pool and condemned the rest of the
        # collection), overallocating threads was doubly risky. Splitting the
        # physical cores across the pool's workers keeps the product
        # (processes x jobs) close to the machine's core count.
        #
        # Floor of 2, not 1: dividing with no floor gives `--jobs 1` in
        # "maximum" mode (physical // workers = 8 // 8), which the C1/
        # parallelism defect review measured as too conservative — it dropped
        # the gain from 6.43x to 5.06x (a 27% drop) on a corpus of 12 3-page
        # PDFs. With a floor of 2, 8 workers x 2 jobs = 16 threads on 8 cores:
        # 2x oversubscription, the usual pattern for a mixed CPU/IO workload,
        # far from the 64 tesseracts of the original behavior.
        ocrmypdf_jobs = max(2, _physical_cores() // workers)
        _convert_in_parallel(conn, config, rows, workers, ocrmypdf_jobs, should_stop, result, language=language)

    record_event(
        conn,
        "conversion",
        "info",
        "log.conversion.summary",
        {
            "converted": result.converted, "ocr": result.ocr_applied,
            "sliced": result.sliced, "blocks": result.blocks_generated, "failed": result.failed,
        },
        language=language,
    )

    return result
